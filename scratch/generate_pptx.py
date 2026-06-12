from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = Presentation()

title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]

slides_data = [
    {
        "type": "title",
        "title": "Predictive HR Analytics: Early Warning System for Infant Attrition",
        "subtitle": "6CS007 Project Presentation\nRidmi Weerasinghe\nStudent ID: 2605498"
    },
    {
        "type": "content",
        "title": "The Challenge in the Garment Sector",
        "bullets": [
            "Sri Lankan apparel industry suffers from severe 'Infant Attrition' (employees quitting within 0-12 months).",
            "Current HR systems only track resignations after they happen.",
            "Lack of proactive, data-driven early warning systems to identify flight risks during the onboarding phase."
        ]
    },
    {
        "type": "content",
        "title": "Aims & Objectives",
        "bullets": [
            "Goal: To develop an Early Warning System using predictive machine learning.",
            "Prove that historical demographic and logistical data can predict turnover.",
            "Build an accessible, secure UI for HR personnel.",
            "Implement 'Explainable AI' so HR managers understand why a risk score is high."
        ]
    },
    {
        "type": "content",
        "title": "Data Engineering: The 12-Month Milestone",
        "bullets": [
            "Raw Dataset: 5,026 employee records (1,500 Active vs 3,526 Inactive).",
            "The Bias Challenge: Standard 'Active vs Inactive' metrics are flawed. An employee who quits after 5 years is a success, not a flight risk.",
            "The Solution: We redefined the target variable specifically around 'Infant Attrition' (0-12 months).",
            "Success (Target 0): Reached 1 year of tenure (2,462 records).",
            "Failure (Target 1): Resigned before 1 year (2,464 records).",
            "Natural Balance: This strategic filtering naturally balanced the dataset at exactly 50/50."
        ]
    },
    {
        "type": "content",
        "title": "System Architecture",
        "bullets": [
            "Frontend (Presentation): Streamlit framework providing a dynamic, Python-based web UI.",
            "Backend (Logic): Scikit-Learn executing the Random Forest Classifier and SHAP calculations.",
            "Database (Storage): MySQL securing User Authentication, RBAC (Role-Based Access Control), and historical batch metrics."
        ]
    },
    {
        "type": "content",
        "title": "Baseline Model Evaluation",
        "bullets": [
            "Tested multiple algorithms: Logistic Regression, Decision Trees, Random Forest, and Gradient Boosting.",
            "Gradient Boosting showed initial statistical promise, but Random Forest was ultimately selected for superior stability against real-world human data.",
            "The Accuracy Paradox: The baseline models prioritized generic 'Accuracy', which often fails to identify the critical minority class (flight risks)."
        ]
    },
    {
        "type": "content",
        "title": "Hyperparameter Tuning & Cross-Validation",
        "bullets": [
            "Execution: Utilized 5-Fold Stratified Cross-Validation to ensure the model generalized well across unseen data.",
            "Optimization Target: Used RandomizedSearchCV, strictly overriding default metrics with 'scoring=recall'.",
            "Optimal Architecture Extracted: The grid search determined the most robust architecture was 500 parallel decision trees (n_estimators=500) restricted to a maximum depth of 10."
        ]
    },
    {
        "type": "content",
        "title": "Final Model Performance Metrics",
        "bullets": [
            "Recall (Sensitivity) -> 98.62%: The most critical metric. The system successfully flags nearly 99% of all true flight risks.",
            "Precision -> 73.57%: When the model flags a risk, it is correct 73.5% of the time (An acceptable false-alarm rate for HR check-ins).",
            "Accuracy -> 81.59%: Strong general correctness across the entire workforce.",
            "F1-Score -> 84.27%: Proves the model maintained mathematical stability despite being heavily biased toward Recall."
        ]
    },
    {
        "type": "content",
        "title": "Explainable AI (SHAP)",
        "bullets": [
            "Predictive models are often 'Black Boxes,' which HR managers distrust.",
            "Integrated SHAP (SHapley Additive exPlanations) to break open the Black Box.",
            "System generates dynamic charts showing the exact Top 5 driving factors behind an employee's risk score (e.g., Age pulling risk down, Job Level pushing risk up)."
        ]
    },
    {
        "type": "content",
        "title": "Product Testing & Quality Assurance",
        "bullets": [
            "Data Validation Testing: Implemented strict schema validation to block malformed CSV uploads before model execution.",
            "Security Testing: Verified SHA-256 password hashing and strict Role-Based Access Control (Admin vs Base User).",
            "Functional UI Testing: Verified that edge cases (like entering an age of '0' or missing columns) trigger clean UI error messages instead of system crashes."
        ]
    },
    {
        "type": "content",
        "title": "System Demonstration: Single Profile",
        "bullets": [
            "(Insert Demo Screenshot Here)",
            "Walkthrough: Demonstrate how a HR user inputs data for a single new hire.",
            "Show how the Explainable AI (SHAP chart) visualizes the risk score instantly."
        ]
    },
    {
        "type": "content",
        "title": "System Demonstration: Batch Analytics",
        "bullets": [
            "(Insert Demo Screenshot Here)",
            "Walkthrough: Demonstrate uploading a CSV for an entire factory batch.",
            "Show the strict schema validation in action.",
            "Display the Executive Dashboard rendering real-time aggregated metrics."
        ]
    },
    {
        "type": "content",
        "title": "Project Conclusion",
        "bullets": [
            "Successfully bridged the gap between theoretical Machine Learning and practical Software Engineering.",
            "Provided a functional PoC (Proof of Concept) that proactively flags risk rather than reacting to it.",
            "Demonstrates how privacy-conscious machine learning can solve real-world operational challenges in Sri Lanka."
        ]
    },
    {
        "type": "title",
        "title": "Thank You!",
        "subtitle": "Any Questions?"
    }
]

for slide_data in slides_data:
    if slide_data["type"] == "title":
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = slide_data["title"]
        subtitle.text = slide_data["subtitle"]
        
    elif slide_data["type"] == "content":
        slide = prs.slides.add_slide(content_slide_layout)
        title = slide.shapes.title
        body_shape = slide.placeholders[1]
        
        title.text = slide_data["title"]
        
        tf = body_shape.text_frame
        tf.clear()
        
        for i, bullet in enumerate(slide_data["bullets"]):
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            if "Target 0" in bullet or "Target 1" in bullet:
                p.level = 1

prs.save('6CS007_Presentation_Ridmi.pptx')
print("Updated Presentation successfully created!")
